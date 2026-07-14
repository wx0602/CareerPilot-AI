import tempfile
import unittest
from pathlib import Path

from models import ContextChunk

from knowledge.parsers import parse_material
from knowledge.rag.service import KnowledgeService, split_text


class ParserTests(unittest.TestCase):
    def test_resume_and_jd_structured_fields(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            resume = root / "resume.txt"
            resume.write_text(
                "张三\n邮箱：zhangsan@example.com\n电话：13800138000\n技能\nJava、MySQL、Docker",
                encoding="utf-8",
            )
            parsed = parse_material(resume, material_type="resume", user_id="u1")
            self.assertIn("Java", parsed.structured_data["skills"])
            self.assertNotIn("SQL", parsed.structured_data["skills"])
            self.assertEqual(parsed.structured_data["phones"], ["13800138000"])

            jd = root / "jd.txt"
            jd.write_text("岗位名称：Java开发\n任职要求\n熟悉Java和MySQL，本科及以上学历。", encoding="utf-8")
            parsed_jd = parse_material(jd, material_type="jd", user_id="u1")
            self.assertEqual(parsed_jd.structured_data["job_title"], "Java开发")
            self.assertIn("MySQL", parsed_jd.structured_data["required_skills"])

    def test_material_index_falls_back_without_chroma(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "project.txt"
            source.write_text("项目使用FastAPI和SQLite，实现简历解析与题库检索。", encoding="utf-8")
            material = parse_material(source, material_type="project_intro", user_id="u1")
            service = KnowledgeService(
                database_path=root / "questions.sqlite3",
                material_store_path=root / "materials.jsonl",
            )
            service.index_material(material)
            result = service.search_materials("简历解析", user_id="u1")
            self.assertEqual(result["total"], 1)
            self.assertEqual(result["contexts"][0]["chunk_id"], material.material_id)
            ContextChunk.model_validate(result["contexts"][0])

    def test_docx_business_plan_and_pptx_pitch_deck(self):
        from docx import Document
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            plan_path = root / "plan.docx"
            document = Document()
            document.add_heading("商业计划书", level=1)
            document.add_paragraph("产品面向高校求职训练场景，采用订阅制服务。")
            document.save(plan_path)
            plan = parse_material(plan_path, material_type="business_plan", user_id="team1")
            self.assertIn("订阅制", plan.parsed_text)

            deck_path = root / "pitch.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "CareerPilot AI"
            slide.placeholders[1].text = "智能求职与创业训练平台"
            presentation.save(deck_path)
            deck = parse_material(deck_path, material_type="pitch_ppt", user_id="team1")
            self.assertIn("智能求职", deck.parsed_text)
            self.assertEqual(deck.type, "pitch_ppt")
            self.assertEqual(deck.structured_data["slide_count"], 1)

    def test_split_text_validates_overlap(self):
        self.assertGreaterEqual(len(split_text("第一段\n" + "内容" * 50, chunk_size=30, overlap=5)), 2)
        with self.assertRaises(ValueError):
            split_text("text", chunk_size=10, overlap=10)


if __name__ == "__main__":
    unittest.main()
