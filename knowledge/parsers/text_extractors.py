from __future__ import annotations

import json
from pathlib import Path
from typing import Callable


class UnsupportedDocumentError(ValueError):
    pass


def _read_plain_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise UnicodeDecodeError("unknown", b"", 0, 1, f"无法识别文件编码: {path.name}")


def _read_json(path: Path) -> str:
    payload = json.loads(_read_plain_text(path))
    return json.dumps(payload, ensure_ascii=False, indent=2)


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("解析 DOCX 需要安装 python-docx") from exc
    document = Document(path)
    blocks = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            values = [" ".join(cell.text.split()) for cell in row.cells]
            if any(values):
                blocks.append(" | ".join(values))
    return "\n".join(blocks)


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("解析 PDF 需要安装 pypdf") from exc
    reader = PdfReader(str(path))
    return "\n".join((page.extract_text() or "").strip() for page in reader.pages).strip()


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("解析 PPTX 需要安装 python-pptx") from exc
    presentation = Presentation(str(path))
    slides: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = " ".join(shape.text.split())
                if text:
                    texts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [" ".join(cell.text.split()) for cell in row.cells]
                    if any(values):
                        texts.append(" | ".join(values))
        if texts:
            slides.append(f"[第{number}页]\n" + "\n".join(texts))
    return "\n\n".join(slides)


EXTRACTORS: dict[str, Callable[[Path], str]] = {
    ".txt": _read_plain_text,
    ".md": _read_plain_text,
    ".csv": _read_plain_text,
    ".json": _read_json,
    ".docx": _read_docx,
    ".pdf": _read_pdf,
    ".pptx": _read_pptx,
}


def extract_text(path: str | Path) -> str:
    source = Path(path)
    if not source.is_file():
        raise FileNotFoundError(f"文件不存在: {source}")
    extractor = EXTRACTORS.get(source.suffix.lower())
    if extractor is None:
        raise UnsupportedDocumentError(
            f"不支持 {source.suffix or '无扩展名'} 文件，支持: {', '.join(sorted(EXTRACTORS))}"
        )
    text = extractor(source)
    normalized = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not normalized:
        raise ValueError(f"未从 {source.name} 提取到文本；扫描版 PDF 需要先 OCR")
    return normalized
