from pathlib import Path
from io import BytesIO

from starlette.testclient import TestClient


def _session(client: TestClient, headers: dict[str, str]) -> str:
    return client.post(
        "/api/training-sessions",
        headers=headers,
        json={"mode": "technical"},
    ).json()["session_id"]


def test_valid_text_upload_is_safely_stored(logged_in, settings):
    client, headers = logged_in
    session_id = _session(client, headers)
    response = client.post(
        "/api/materials/upload",
        headers=headers,
        data={"session_id": session_id, "material_type": "resume"},
        files={"file": ("../../resume.txt", "后端开发经历", "text/plain")},
    )
    assert response.status_code == 201
    body = response.json()
    assert body["filename"] == "resume.txt"
    assert body["parse_status"] == "parsed"
    stored = list(Path(settings.upload_dir).glob("*.txt"))
    assert len(stored) == 1
    assert stored[0].parent.resolve() == Path(settings.upload_dir).resolve()


def test_rejects_forged_or_oversized_file(logged_in):
    client, headers = logged_in
    session_id = _session(client, headers)
    forged = client.post(
        "/api/materials/upload",
        headers=headers,
        data={"session_id": session_id, "material_type": "resume"},
        files={"file": ("resume.pdf", b"not a pdf", "application/pdf")},
    )
    assert forged.status_code == 415
    assert forged.json()["detail"]["code"] == "invalid_file_content"

    oversized = client.post(
        "/api/materials/upload",
        headers=headers,
        data={"session_id": session_id, "material_type": "jd"},
        files={"file": ("jd.txt", b"a" * 1025, "text/plain")},
    )
    assert oversized.status_code == 413


def test_rejects_phase_two_material_type(logged_in):
    client, headers = logged_in
    session_id = _session(client, headers)
    response = client.post(
        "/api/materials/upload",
        headers=headers,
        data={"session_id": session_id, "material_type": "pitch_deck"},
        files={"file": ("plan.txt", "test", "text/plain")},
    )
    assert response.status_code == 422


def test_accepts_pitch_pptx_material(logged_in):
    from pptx import Presentation

    client, headers = logged_in
    client.app.state.settings.max_upload_bytes = 1024 * 1024
    session_id = _session(client, headers)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "CareerPilot 路演"
    slide.placeholders[1].text = "目标用户与商业模式"
    content = BytesIO()
    presentation.save(content)

    response = client.post(
        "/api/materials/upload",
        headers=headers,
        data={"session_id": session_id, "material_type": "pitch_ppt"},
        files={
            "file": (
                "pitch.pptx",
                content.getvalue(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert response.status_code == 201
    assert response.json()["material_type"] == "pitch_ppt"
    assert response.json()["parse_status"] == "parsed"
